from pptx import Presentation
prs = Presentation("examples/sample_project/project_intro_deck.pptx")
print(f"Total slides: {len(prs.slides)}")
for i, s in enumerate(prs.slides, 1):
    texts = [sh.text_frame.text[:55].replace("\n"," ") for sh in s.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
    pics  = sum(1 for sh in s.shapes if sh.shape_type == 13)
    label = texts[0] if texts else "(empty)"
    img   = " [image]" if pics else ""
    print(f"  Slide {i:2d}: {label}{img}")
