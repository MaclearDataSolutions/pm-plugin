"""
create_project_intro_deck.py
Template-driven slide deck generator for both intro and update decks.

Usage:
  python tools/scripts/create_project_intro_deck.py --project project --type intro
  python tools/scripts/create_project_intro_deck.py --project project --type update
  python tools/scripts/create_project_intro_deck.py --project project --config tools/templates/slide_template_intro.json
"""

import re, sys, json, csv, argparse
from pathlib import Path
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

TOOLS_DIR     = Path(__file__).parent           # tools/scripts/
TEMPLATES_DIR = TOOLS_DIR.parent / "templates"  # tools/templates/
ROOT          = TOOLS_DIR.parent.parent          # repo root

DEFAULT_CONFIGS = {
    "intro":  TEMPLATES_DIR / "slide_template_intro.json",
    "update": TEMPLATES_DIR / "slide_template_update.json",
}

# ---------------------------------------------------------------------------
# Colors
# ---------------------------------------------------------------------------
def rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

DARK_BLUE  = rgb("1F3864")
MID_BLUE   = rgb("2E75B6")
LIGHT_BLUE = rgb("DEEAF1")
AMBER      = rgb("FFC000")
WHITE      = rgb("FFFFFF")
DARK_TEXT  = rgb("1A1A1A")
LIGHT_GREY = rgb("F2F2F2")
BROWN      = rgb("843C0C")
GREEN      = rgb("70AD47")

LAYOUT_MAP = {
    "Title Slide": 0, "Title and Content": 1,
    "Section Header": 2, "Title Only": 5, "Blank": 6,
}

# ---------------------------------------------------------------------------
# Config + parsing
# ---------------------------------------------------------------------------
def load_config(path: Path) -> dict:
    if not path.exists():
        sys.exit(f"ERROR: Config not found: {path}")
    with path.open(encoding="utf-8") as f:
        return json.load(f)

def parse_project_plan(path: Path) -> dict:
    if not path.exists():
        sys.exit(f"ERROR: project_plan.md not found: {path}")
    text  = path.read_text(encoding="utf-8")
    parts = re.split(r"\n## ", text)
    secs  = {}
    for part in parts[1:]:
        lines = part.strip().splitlines()
        secs[lines[0].strip()] = "\n".join(lines[1:]).strip()
    return secs

def find_section(secs: dict, keyword: str) -> str:
    kw = keyword.lower()
    for h, c in secs.items():
        if kw in h.lower():
            return c
    return ""

def parse_csv_file(path: Path) -> list:
    if not path.exists():
        return []
    with path.open(encoding="utf-8") as f:
        return list(csv.DictReader(f))

def parse_update_notes(path: Path) -> list:
    """Extract bullet points from update.md if it exists."""
    if not path.exists():
        return []
    text = path.read_text(encoding="utf-8")
    bullets = []
    for line in text.splitlines():
        line = line.strip()
        if line.startswith(("- ", "* ")):
            bullets.append(line[2:].strip())
        elif re.match(r"^\d+\.\s", line):
            bullets.append(re.sub(r"^\d+\.\s+", "", line))
    return bullets

def extract_bullets(content: str, max_n: int = 99) -> list:
    out = []
    for line in content.splitlines():
        line = line.strip()
        if line.startswith(("- ", "* ")):
            t = line[2:].strip()
            if t: out.append(t)
        elif re.match(r"^\d+\.\s", line):
            out.append(re.sub(r"^\d+\.\s+", "", line))
        if len(out) >= max_n: break
    return out

def extract_subsection(content: str, name: str) -> str:
    m = re.search(rf"###\s+{re.escape(name)}\s*\n(.*?)(?=\n###|\Z)", content, re.DOTALL|re.IGNORECASE)
    return m.group(1).strip() if m else ""

def extract_table(content: str) -> tuple:
    lines = [l.strip() for l in content.splitlines() if l.strip().startswith("|")]
    if len(lines) < 3: return [], []
    headers = [h.strip() for h in lines[0].strip("|").split("|")]
    rows = []
    for line in lines[2:]:
        cells = [c.strip() for c in line.strip("|").split("|")]
        while len(cells) < len(headers): cells.append("")
        rows.append(dict(zip(headers, cells)))
    return headers, rows

def get_plan_owner(secs: dict) -> str:
    gov = find_section(secs, "3. Project identity")
    for line in gov.splitlines():
        if "day-to-day contact" in line.lower() or "contact" in line.lower():
            val = line.split(":", 1)[-1].strip().lstrip("-").strip()
            if val: return val
    return "Project Owner"

# ---------------------------------------------------------------------------
# Gantt image export
# ---------------------------------------------------------------------------
def export_gantt_image(excel_path: Path, sheet_name: str, output_path: Path) -> bool:
    try:
        import win32com.client
        from PIL import ImageGrab
    except ImportError as e:
        print(f"  WARNING: Cannot export Gantt image — {e}"); return False
    if not excel_path.exists():
        print(f"  WARNING: Excel not found: {excel_path}"); return False
    print(f"  Opening {excel_path.name} / '{sheet_name}' ...")
    xl = win32com.client.Dispatch("Excel.Application")
    xl.Visible = True
    xl.DisplayAlerts = False
    try:
        wb = xl.Workbooks.Open(str(excel_path.absolute()))
        ws = wb.Sheets(sheet_name)
        ws.UsedRange.CopyPicture(Appearance=1, Format=2)
        img = ImageGrab.grabclipboard()
        if img is None: print("  WARNING: Clipboard empty."); return False
        img.save(str(output_path), "PNG")
        print(f"  Saved: {output_path.name}"); return True
    except Exception as e:
        print(f"  WARNING: Export failed — {e}"); return False
    finally:
        try: wb.Close(False); xl.Quit()
        except: pass

# ---------------------------------------------------------------------------
# Slide shape helpers
# ---------------------------------------------------------------------------
def clear_placeholders(slide):
    for ph in slide.placeholders:
        ph._element.getparent().remove(ph._element)

def font_name(config: dict) -> str:
    return config.get("style", {}).get("font_title", "Aptos Display")

def body_font(config: dict) -> str:
    return config.get("style", {}).get("font_body", "Aptos")

def primary(config: dict) -> RGBColor:
    return rgb(config.get("style", {}).get("primary_color", "1F3864"))

def accent(config: dict) -> RGBColor:
    return rgb(config.get("style", {}).get("accent_color", "FFC000"))

def add_title_bar(slide, prs, title_text: str, config: dict):
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = primary(config); bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = False
    tf.margin_left = Inches(0.32); tf.margin_top = Inches(0.26)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title_text
    r.font.name = font_name(config); r.font.bold = True
    r.font.size = Pt(28); r.font.color.rgb = WHITE

def add_descriptor(slide, prs, text: str, config: dict):
    tb = slide.shapes.add_textbox(Inches(0.32), Inches(1.07), prs.slide_width - Inches(0.64), Inches(0.35))
    tf = tb.text_frame; p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.name = body_font(config); r.font.size = Pt(13); r.font.color.rgb = MID_BLUE

def add_separator(slide, prs, config: dict):
    line = slide.shapes.add_shape(1, Inches(0.32), Inches(1.44), prs.slide_width - Inches(0.64), Pt(1.5))
    line.fill.solid(); line.fill.fore_color.rgb = accent(config); line.line.fill.background()

def add_footer(slide, prs, text: str, config: dict):
    W, H = prs.slide_width, prs.slide_height
    tb = slide.shapes.add_textbox(Inches(0.32), H - Inches(0.38), W - Inches(0.64), Inches(0.3))
    tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text
    r.font.name = body_font(config); r.font.size = Pt(9); r.font.color.rgb = rgb("808080")

def content_area(prs):
    return Inches(0.32), Inches(1.55), prs.slide_width - Inches(0.64), prs.slide_height - Inches(2.0)

def resolve_footer(slide_def: dict, config: dict, today: str) -> str:
    if not slide_def.get("footer"): return ""
    tmpl = config.get("style",{}).get("footer_text", "{company} | {project_name} | Confidential")
    return tmpl.replace("{project_name}", config["meta"].get("project_name","Project")).replace("{date}", today)

def set_cell_bg_color(cell, color: RGBColor):
    tc = cell._tc; tcPr = tc.get_or_add_tcPr()
    old = tcPr.find(qn("a:solidFill"))
    if old is not None: tcPr.remove(old)
    fill = parse_xml(
        f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:srgbClr val="{str(color):>06s}"/></a:solidFill>')
    tcPr.insert(0, fill)

def fmt_cell(cell, text, bold=False, color=None, size=10, align=PP_ALIGN.LEFT, font="Aptos"):
    color = color or DARK_TEXT
    cell.text = str(text); tf = cell.text_frame; tf.word_wrap = True
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.name = font; run.font.bold = bold
            run.font.size = Pt(size); run.font.color.rgb = color

def build_table_shape(slide, prs, col_headers, display_rows, col_map,
                      config, L, T, W, H, max_rows=None):
    if max_rows: display_rows = display_rows[:max_rows]
    if not col_headers or not display_rows:
        tb = slide.shapes.add_textbox(L, T, W, H)
        tb.text_frame.paragraphs[0].add_run().text = "No data found in project plan."
        return
    n_rows = len(display_rows) + 1; n_cols = len(col_headers)
    tbl = slide.shapes.add_table(n_rows, n_cols, L, T, W, H).table
    col_w = W // n_cols
    for i in range(n_cols): tbl.columns[i].width = col_w
    bf = body_font(config)
    for c, ch in enumerate(col_headers):
        cell = tbl.cell(0, c)
        set_cell_bg_color(cell, primary(config))
        fmt_cell(cell, ch, bold=True, color=WHITE, size=11, align=PP_ALIGN.CENTER, font=bf)
    for r, row in enumerate(display_rows, 1):
        bg = LIGHT_BLUE if r % 2 == 0 else WHITE
        for c, ch in enumerate(col_headers):
            cell = tbl.cell(r, c)
            set_cell_bg_color(cell, bg)
            fmt_cell(cell, row.get(col_map.get(ch, ch), ""), size=10, font=bf)

def map_columns(col_headers, md_headers):
    col_map = {}
    for ch in col_headers:
        for mh in md_headers:
            if ch.lower() in mh.lower() or mh.lower() in ch.lower():
                col_map[ch] = mh; break
        if ch not in col_map: col_map[ch] = ch
    return col_map

# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def build_title_slide(prs, sd, config, secs, today, extra):
    layout = prs.slide_layouts[LAYOUT_MAP.get(sd.get("layout","Blank"), 6)]
    slide  = prs.slides.add_slide(layout); clear_placeholders(slide)
    W, H = prs.slide_width, prs.slide_height
    p_col = primary(config); a_col = accent(config)

    bg = slide.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = p_col; bg.line.fill.background()
    strip = slide.shapes.add_shape(1, 0, H - Inches(0.28), W, Inches(0.28))
    strip.fill.solid(); strip.fill.fore_color.rgb = a_col; strip.line.fill.background()
    bar = slide.shapes.add_shape(1, 0, Inches(1.9), Inches(0.07), Inches(3.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = a_col; bar.line.fill.background()

    fn = font_name(config); bn = body_font(config)
    def add_tb(left, top, w, h, text, size, bold=False, color=WHITE):
        tb = slide.shapes.add_textbox(left, top, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = text
        r.font.name = fn if size >= 20 else bn
        r.font.bold = bold; r.font.size = Pt(size); r.font.color.rgb = color

    add_tb(Inches(0.7), Inches(2.0), W - Inches(1.4), Inches(1.8),
           sd.get("title","").replace("{project_name}", config["meta"].get("project_name","Project")),
           40, bold=True)
    add_tb(Inches(0.7), Inches(3.8), W - Inches(1.4), Inches(0.8),
           sd.get("subtitle","Project Introduction"), 24, color=LIGHT_BLUE)
    owner = extra.get("owner","Project Owner")
    author = sd.get("author_line","{owner} | {company} | {date}")\
        .replace("{owner}", owner).replace("{date}", today)
    add_tb(Inches(0.7), Inches(6.5), W - Inches(1.4), Inches(0.5), author, 13, color=LIGHT_BLUE)
    conf = slide.shapes.add_textbox(Inches(0.7), H - Inches(0.5), W - Inches(1.4), Inches(0.3))
    cp = conf.text_frame.paragraphs[0]; cp.alignment = PP_ALIGN.LEFT
    cr = cp.add_run(); cr.text = "Confidential - Internal Use"
    cr.font.name = bn; cr.font.size = Pt(10); cr.font.color.rgb = LIGHT_BLUE

    if sd.get("speaker_notes"):
        slide.notes_slide.notes_text_frame.text = sd["speaker_notes"]

def build_bullets_slide(prs, sd, config, secs, today, extra):
    layout = prs.slide_layouts[LAYOUT_MAP.get(sd.get("layout","Blank"), 6)]
    slide  = prs.slides.add_slide(layout); clear_placeholders(slide)
    add_title_bar(slide, prs, sd.get("title",""), config)
    if sd.get("descriptor"): add_descriptor(slide, prs, sd["descriptor"], config)
    add_separator(slide, prs, config)
    footer_text = resolve_footer(sd, config, today)
    if footer_text: add_footer(slide, prs, footer_text, config)

    src = sd.get("content_source", {})
    content = find_section(secs, src.get("section",""))
    fields  = src.get("fields", [])
    prefix  = sd.get("bullet_prefix", "• ")
    max_b   = sd.get("max_bullets", 6)
    bn      = body_font(config)

    bullets = []
    if fields:
        for field in fields:
            sub = extract_subsection(content, field)
            if sub:
                items = extract_bullets(sub) or [s.strip().rstrip(".") for s in re.split(r"\.\s+",sub) if s.strip()]
                bullets.extend(items)
            else:
                for line in content.splitlines():
                    if field.lower() in line.lower() and ":" in line:
                        val = line.split(":",1)[-1].strip().lstrip("-").strip()
                        if val: bullets.append(val)
    else:
        bullets = extract_bullets(content)
    bullets = [b for b in bullets if b][:max_b]

    L, T, W, H = content_area(prs)
    tb = slide.shapes.add_textbox(L, T, W, H)
    tf = tb.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.05)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(7); p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = f"{prefix}{b}"
        r.font.name = bn; r.font.size = Pt(sd.get("font_size", 15)); r.font.color.rgb = DARK_TEXT
    if not bullets:
        tf.paragraphs[0].add_run().text = "Content not found."

    if sd.get("speaker_notes"):
        slide.notes_slide.notes_text_frame.text = sd["speaker_notes"]

def build_two_column_slide(prs, sd, config, secs, today, extra):
    layout = prs.slide_layouts[LAYOUT_MAP.get(sd.get("layout","Blank"), 6)]
    slide  = prs.slides.add_slide(layout); clear_placeholders(slide)
    add_title_bar(slide, prs, sd.get("title",""), config)
    if sd.get("descriptor"): add_descriptor(slide, prs, sd["descriptor"], config)
    add_separator(slide, prs, config)
    footer_text = resolve_footer(sd, config, today)
    if footer_text: add_footer(slide, prs, footer_text, config)

    W, H = prs.slide_width, prs.slide_height
    col_defs = sd.get("columns", [])
    n_cols = len(col_defs) or 2
    pad = Inches(0.32); gap = Inches(0.2)
    avail_w = W - 2 * pad - (n_cols - 1) * gap
    col_w = avail_w // n_cols
    top = Inches(1.58); bot = H - Inches(0.55); col_h = bot - top
    bn = body_font(config); fn = font_name(config)

    for idx, col_def in enumerate(col_defs):
        left = pad + idx * (col_w + gap)
        hdr_c = rgb(col_def.get("header_color", "1F3864"))

        # Column header bar
        hbar = slide.shapes.add_shape(1, left, top, col_w, Inches(0.38))
        hbar.fill.solid(); hbar.fill.fore_color.rgb = hdr_c; hbar.line.fill.background()
        htf = hbar.text_frame; htf.margin_left = Inches(0.12); htf.margin_top = Inches(0.06)
        hp = htf.paragraphs[0]; hr = hp.add_run(); hr.text = col_def.get("header","")
        hr.font.name = fn; hr.font.bold = True; hr.font.size = Pt(13); hr.font.color.rgb = WHITE

        # Bullet content
        src = col_def.get("content_source", {})
        content = find_section(secs, src.get("section",""))
        fields  = src.get("fields", [])
        max_b   = col_def.get("max_bullets", 8)
        bullets = []
        for field in fields:
            sub = extract_subsection(content, field)
            if sub: bullets.extend(extract_bullets(sub))
            else: bullets.extend(extract_bullets(content))
        bullets = [b for b in bullets if b][:max_b]

        btb = slide.shapes.add_textbox(left + Inches(0.08), top + Inches(0.42),
                                       col_w - Inches(0.16), col_h - Inches(0.44))
        btf = btb.text_frame; btf.word_wrap = True; btf.margin_left = Inches(0.05)
        for i, b in enumerate(bullets):
            p2 = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            p2.space_before = Pt(5); r2 = p2.add_run()
            r2.text = f"•  {b}"; r2.font.name = bn
            r2.font.size = Pt(12); r2.font.color.rgb = DARK_TEXT

    if sd.get("speaker_notes"):
        slide.notes_slide.notes_text_frame.text = sd["speaker_notes"]

def build_table_slide(prs, sd, config, secs, today, extra):
    layout = prs.slide_layouts[LAYOUT_MAP.get(sd.get("layout","Blank"), 6)]
    slide  = prs.slides.add_slide(layout); clear_placeholders(slide)
    add_title_bar(slide, prs, sd.get("title",""), config)
    descriptor = sd.get("descriptor","").replace("{date}", today)
    if descriptor: add_descriptor(slide, prs, descriptor, config)
    add_separator(slide, prs, config)
    footer_text = resolve_footer(sd, config, today)
    if footer_text: add_footer(slide, prs, footer_text, config)

    src         = sd.get("content_source", {})
    content     = find_section(secs, src.get("section",""))
    col_headers = sd.get("table_columns", [])
    max_rows    = sd.get("max_rows", 5)
    md_headers, md_rows = extract_table(content)
    if not md_rows:
        for field in src.get("fields",[]):
            sub = extract_subsection(content, field)
            if sub: md_headers, md_rows = extract_table(sub)
            if md_rows: break
    if col_headers and md_headers: col_map = map_columns(col_headers, md_headers)
    else: col_headers = md_headers[:6]; col_map = {h: h for h in col_headers}

    L, T, W, H = content_area(prs)
    build_table_shape(slide, prs, col_headers, md_rows, col_map, config, L, T, W, H, max_rows)

    if sd.get("speaker_notes"):
        slide.notes_slide.notes_text_frame.text = sd["speaker_notes"]

def build_milestone_status_slide(prs, sd, config, csv_rows, today, extra):
    """Reads milestone rows from gantt_tasks.csv and shows progress table."""
    layout = prs.slide_layouts[LAYOUT_MAP.get(sd.get("layout","Blank"), 6)]
    slide  = prs.slides.add_slide(layout); clear_placeholders(slide)
    add_title_bar(slide, prs, sd.get("title",""), config)
    descriptor = sd.get("descriptor","").replace("{date}", today)
    if descriptor: add_descriptor(slide, prs, descriptor, config)
    add_separator(slide, prs, config)
    footer_text = resolve_footer(sd, config, today)
    if footer_text: add_footer(slide, prs, footer_text, config)

    filt = sd.get("content_source",{}).get("filter",{})
    filt_col = filt.get("column","Milestone?"); filt_val = filt.get("value","Yes").lower()
    ms_rows = [r for r in csv_rows if r.get(filt_col,"").strip().lower() == filt_val]

    col_headers = sd.get("table_columns", ["Task ID","Task","Milestone date","Owner","Progress","Status"])
    col_map     = {ch: ch for ch in col_headers}
    L, T, W, H  = content_area(prs)
    build_table_shape(slide, prs, col_headers, ms_rows, col_map, config, L, T, W, H, sd.get("max_rows",6))

    if sd.get("speaker_notes"):
        slide.notes_slide.notes_text_frame.text = sd["speaker_notes"]

def build_update_summary_slide(prs, sd, config, secs, update_bullets, today, extra):
    """Executive Summary slide for the update deck — reads from update.md or fallback."""
    layout = prs.slide_layouts[LAYOUT_MAP.get(sd.get("layout","Blank"), 6)]
    slide  = prs.slides.add_slide(layout); clear_placeholders(slide)
    add_title_bar(slide, prs, sd.get("title","Executive Summary"), config)
    if sd.get("descriptor"): add_descriptor(slide, prs, sd["descriptor"], config)
    add_separator(slide, prs, config)
    footer_text = resolve_footer(sd, config, today)
    if footer_text: add_footer(slide, prs, footer_text, config)

    src = sd.get("content_source",{})
    bullets = update_bullets if update_bullets else []
    if not bullets:
        fb_sec = src.get("fallback_section","21. Recommended next steps")
        content = find_section(secs, fb_sec)
        bullets = extract_bullets(content)
    bullets = [b for b in bullets if b][:sd.get("max_bullets",6)]

    bn = body_font(config)
    L, T, W, H = content_area(prs)
    tb = slide.shapes.add_textbox(L, T, W, H)
    tf = tb.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(7); r = p.add_run()
        r.text = f"•  {b}"; r.font.name = bn
        r.font.size = Pt(15); r.font.color.rgb = DARK_TEXT
    if not bullets:
        tf.paragraphs[0].add_run().text = "No update notes found. Add content to update.md."

    if sd.get("speaker_notes"):
        slide.notes_slide.notes_text_frame.text = sd["speaker_notes"]

def build_gantt_slide(prs, sd, config, img_path, today, extra):
    layout = prs.slide_layouts[LAYOUT_MAP.get(sd.get("layout","Blank"), 6)]
    slide  = prs.slides.add_slide(layout); clear_placeholders(slide)
    add_title_bar(slide, prs, sd.get("title","Project Timeline"), config)
    if sd.get("descriptor"): add_descriptor(slide, prs, sd["descriptor"], config)
    add_separator(slide, prs, config)
    footer_text = resolve_footer(sd, config, today)
    if footer_text: add_footer(slide, prs, footer_text, config)

    W, H = prs.slide_width, prs.slide_height
    pos  = sd.get("image_position",{"left_pct":0.02,"top_pct":0.16,"width_pct":0.96,"height_pct":0.78})
    left  = int(pos["left_pct"]  * W); top   = int(pos["top_pct"]   * H)
    width = int(pos["width_pct"] * W); height= int(pos["height_pct"]* H)

    if img_path and img_path.exists():
        slide.shapes.add_picture(str(img_path), left, top, width, height)
    else:
        tb = slide.shapes.add_textbox(left, top, width, height)
        r  = tb.text_frame.paragraphs[0].add_run()
        r.text = "Gantt image not available — ensure gantt_chart.xlsx is present."
        r.font.size = Pt(14); r.font.color.rgb = BROWN

    if sd.get("speaker_notes"):
        slide.notes_slide.notes_text_frame.text = sd["speaker_notes"]

# ---------------------------------------------------------------------------
# Builder dispatch
# ---------------------------------------------------------------------------
def build_slide(prs, sd, config, secs, csv_rows, update_bullets, gantt_img, today, extra):
    ctype = sd.get("content_type","bullets")
    kwargs = dict(config=config, today=today, extra=extra)
    if   ctype == "title_only":       build_title_slide(prs, sd, **kwargs, secs=secs)
    elif ctype == "bullets":          build_bullets_slide(prs, sd, **kwargs, secs=secs)
    elif ctype == "two_column_bullets": build_two_column_slide(prs, sd, **kwargs, secs=secs)
    elif ctype == "table":            build_table_slide(prs, sd, **kwargs, secs=secs)
    elif ctype == "gantt_image":      build_gantt_slide(prs, sd, **kwargs, img_path=gantt_img)
    elif ctype == "milestone_status": build_milestone_status_slide(prs, sd, **kwargs, csv_rows=csv_rows)
    elif ctype == "update_summary":   build_update_summary_slide(prs, sd, **kwargs, secs=secs, update_bullets=update_bullets)
    else: print(f"  WARNING: Unknown content_type '{ctype}', skipping.")

# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--project", help="Project folder (e.g. 0608)")
    parser.add_argument("--type",    choices=["intro","update"], default="intro")
    parser.add_argument("--config",  help="Path to slide template JSON")
    parser.add_argument("--output",  help="Output .pptx path")
    args = parser.parse_args()

    # Resolve config
    if args.config:
        config_path = ROOT / args.config
    else:
        config_path = DEFAULT_CONFIGS[args.type]

    print(f"Config:   {config_path}")
    config = load_config(config_path)

    project_dir = ROOT / args.project if args.project else ROOT
    sources     = config.get("sources", {})
    plan_path   = project_dir / sources.get("project_plan","project_plan.md")
    csv_path    = project_dir / sources.get("gantt_csv","gantt_tasks.csv")
    gantt_xlsx  = project_dir / sources.get("gantt_excel","gantt_chart.xlsx")
    update_path = project_dir / sources.get("update_notes","update.md")
    _tmpl_rel = config.get("meta", {}).get("pptx_template")
    pptx_tmpl = (ROOT / _tmpl_rel) if _tmpl_rel else Path()
    if not pptx_tmpl.exists():
        pptx_tmpl = TEMPLATES_DIR / (
            "slide_template_intro.pptx" if args.type == "intro" else "slide_template_update.pptx"
        )
    out_name    = f"project_{'intro' if args.type == 'intro' else 'update'}_deck.pptx"
    out_path    = ROOT / args.output if args.output else project_dir / out_name
    today       = date.today().strftime("%B %d, %Y")

    print(f"Plan:     {plan_path}")
    print(f"Template: {pptx_tmpl}")
    print(f"Output:   {out_path}\n")

    print("Parsing project_plan.md ...")
    secs = parse_project_plan(plan_path)
    print(f"  {len(secs)} sections.")
    owner = get_plan_owner(secs)

    csv_rows      = parse_csv_file(csv_path)
    update_bullets= parse_update_notes(update_path)
    if update_bullets: print(f"  update.md: {len(update_bullets)} bullets.")

    # Export Gantt image (intro deck only)
    gantt_img = None
    if any(sd.get("content_type") == "gantt_image" for sd in config.get("slides",[])):
        gantt_cfg = config.get("gantt_image", {})
        img_path  = project_dir / gantt_cfg.get("export_path","gantt_timeline_export.png")
        print("Exporting Gantt image ...")
        if export_gantt_image(gantt_xlsx, gantt_cfg.get("sheet","Gantt Timeline"), img_path):
            gantt_img = img_path

    prs = Presentation(str(pptx_tmpl)) if pptx_tmpl.exists() else Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Remove all existing slides from the template — keep only the master/theme
    slide_id_list = prs.slides._sldIdLst
    for sld_id in list(slide_id_list):
        slide_id_list.remove(sld_id)

    print("Building slides ...")
    extra = {"owner": owner}
    for sd in config.get("slides", []):
        print(f"  {sd.get('id','?')}: {sd.get('name','')} [{sd.get('content_type','')}]")
        build_slide(prs, sd, config, secs, csv_rows, update_bullets, gantt_img, today, extra)

    prs.save(str(out_path))
    print(f"\nSaved: {out_path}")
    if gantt_img and gantt_img.exists():
        gantt_img.unlink()
    print("Done.")

if __name__ == "__main__":
    main()
