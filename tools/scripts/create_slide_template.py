"""
Creates tools/templates/slide_template_intro.pptx — the base visual template for the slide deck.
Run once when you need to regenerate the base template from scratch.
The generation script (create_project_intro_deck.py) opens this file
and populates it with project data.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path

OUT = Path(__file__).parent.parent / "templates" / "slide_template_intro.pptx"

# ---------------------------------------------------------------------------
# Default color palette
# ---------------------------------------------------------------------------
DARK_BLUE  = "1F3864"
MID_BLUE   = "2E75B6"
LIGHT_BLUE = "DEEAF1"
AMBER      = "FFC000"
WHITE      = "FFFFFF"
DARK_TEXT  = "1A1A1A"
LIGHT_GREY = "F2F2F2"


def set_theme_colors(prs):
    """Replace the default theme colors with custom colors in the XML."""
    theme_el = prs.slide_master.element.find(
        ".//{http://schemas.openxmlformats.org/drawingml/2006/main}theme"
    )
    if theme_el is None:
        return

    # Find the color scheme element and update dk1 (dark1) and accent1
    NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    fmtScheme = theme_el.find(f".//{{{NS}}}fmtScheme")

    # Update dk1 (primary dark) to DARK_BLUE
    for srgbClr in theme_el.findall(f".//{{{NS}}}dk1/{{{NS}}}srgbClr"):
        srgbClr.set("val", DARK_BLUE)
    # Update accent1 to MID_BLUE
    for srgbClr in theme_el.findall(f".//{{{NS}}}accent1/{{{NS}}}srgbClr"):
        srgbClr.set("val", MID_BLUE)
    # Update accent2 to AMBER
    for srgbClr in theme_el.findall(f".//{{{NS}}}accent2/{{{NS}}}srgbClr"):
        srgbClr.set("val", AMBER)


def style_master_title(prs):
    """Set the slide master title placeholder to white bold text."""
    master = prs.slide_master
    for ph in master.placeholders:
        tf = ph.text_frame if ph.has_text_frame else None
        if tf is None:
            continue
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def main():
    prs = Presentation()

    # 16:9 widescreen
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    set_theme_colors(prs)
    style_master_title(prs)

    prs.save(str(OUT))
    print(f"Saved: {OUT}")
    print("Layouts available:")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"  {i}: {layout.name}")


if __name__ == "__main__":
    main()
