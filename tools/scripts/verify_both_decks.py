from pptx import Presentation
from pathlib import Path

for f, label in [("examples/sample_project/project_intro_deck.pptx","INTRO"), ("project/project_update_deck.pptx","UPDATE")]:
    p   = Presentation(f)
    kb  = round(Path(f).stat().st_size / 1024)
    print(f"{label}: {len(p.slides)} slides, {kb}KB")
    for i, s in enumerate(p.slides, 1):
        texts = [sh.text_frame.text[:48].replace("\n"," ") for sh in s.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
        pics  = sum(1 for sh in s.shapes if sh.shape_type == 13)
        tbls  = sum(1 for sh in s.shapes if sh.has_table)
        label2 = texts[0] if texts else "(no text)"
        extras = ("" + (" [img]" if pics else "") + (f" [tbl x{tbls}]" if tbls else ""))
        print(f"  {i:2d}: {label2}{extras}")
    print()
