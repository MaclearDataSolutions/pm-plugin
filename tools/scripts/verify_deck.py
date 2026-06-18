from pptx import Presentation
from pathlib import Path

prs = Presentation("examples/sample_project/project_intro_deck.pptx")
print(f"Slides : {len(prs.slides)}")
print(f"Size   : {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\" (16:9)")
print()
for i, slide in enumerate(prs.slides, 1):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                texts.append(t[:70].replace("\n", " "))
    title = texts[0] if texts else "(no text)"
    print(f"  Slide {i:2d}: {title}")
    notes = slide.notes_slide.notes_text_frame.text.strip()
    if notes:
        preview = notes[:80] + "..." if len(notes) > 80 else notes
        print(f"           Notes: {preview}")
    # Check for images
    pics = [s for s in slide.shapes if s.shape_type == 13]
    if pics:
        print(f"           Images: {len(pics)}")
    # Check for tables
    tbls = [s for s in slide.shapes if s.has_table]
    if tbls:
        rows = tbls[0].table.rows
        print(f"           Table: {len(rows)} rows x {len(rows[0].cells)} cols")
